import streamlit as st
import pandas as pd
import numpy as np
import itertools
import matplotlib.pyplot as plt
import seaborn as sns
from collections import Counter
from scipy.stats import skew
import io
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

# Optional: for GPT (enable only if you have API key)
try:
    import openai
    openai.api_key = st.secrets["openai_key"]  # or hardcode your key
except:
    pass

# ----------------------------
# Streamlit App Configuration
# ----------------------------
st.set_page_config(page_title="TURF Analysis Agent", layout="centered")
st.title("🤖 TURF Analysis Assistant")

# ----------------------------
# Step Control
# ----------------------------
if 'step' not in st.session_state:
    st.session_state.step = 1
if 'df' not in st.session_state:
    st.session_state.df = None

# ----------------------------
# Step 1: Upload Excel File
# ----------------------------
if st.session_state.step == 1:
    st.header("Step 1: Upload PET Message Data")
    uploaded_file = st.file_uploader("Upload your Excel file", type=["xlsx"])
    if uploaded_file:
        st.session_state.df = pd.read_excel(uploaded_file)
        st.success("✅ File loaded successfully.")
        st.session_state.step += 1
        st.rerun()

# ----------------------------
# Step 2: Data Summary + GPT Suggestion
# ----------------------------
elif st.session_state.step == 2:
    st.header("Step 2: Data Summary & AM/GM Recommendation")
    df = st.session_state.df
    score_types = ['Differentiated', 'Believable', 'Motivating']
    message_ids = sorted(set(col.split("_")[0] for col in df.columns if col.startswith("M")))
    st.write(f"📊 Respondents: {df.shape[0]}")
    st.write(f"💬 Messages: {', '.join(message_ids)}")

    # ✅ Prepare summary_df
    summary_rows = []
    for msg in message_ids:
        for score_type in score_types:
            col_name = f"{msg}_{score_type}"
            scores = df[col_name].dropna()
            summary_rows.append({
                "Message": msg,
                "ScoreType": score_type,
                "Mean": round(scores.mean(), 2),
                "StdDev": round(scores.std(), 2),
                "Skew": round(skew(scores), 2)
            })
    summary_df = pd.DataFrame(summary_rows)

    # ✅ Run GPT only once
    if "gpt_recommendation" not in st.session_state:
        try:
            grouped_summary = summary_df.groupby("ScoreType")
            prompt = (
                "You are a message effectiveness analyst. Based on the following summary of "
                "Top-2-Box PET message scores, recommend whether to use Arithmetic Mean or "
                "Geometric Mean to combine Differentiated, Believable, and Motivating scores "
                "for each message.\n\n"
            )

            for name, group in grouped_summary:
                prompt += f"\n--- {name.upper()} ---\n"
                for _, row in group.iterrows():
                    prompt += (
                        f"{row['Message']}: Mean={row['Mean']}, "
                        f"StdDev={row['StdDev']}, Skew={row['Skew']}\n"
                    )

            prompt += (
                "\nPlease recommend whether Arithmetic Mean or Geometric Mean is better and why, "
                "in 2-3 sentences."
            )

            key = st.secrets["openai_key"]
            client = openai.OpenAI(api_key=key)
            response = client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are an expert in pharmaceutical message testing and analytics."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1
            )
            st.session_state["gpt_recommendation"] = response.choices[0].message.content
        except Exception as e:
            st.session_state["gpt_recommendation"] = f"⚠️ GPT recommendation failed: {e}"

    # ✅ Always display the recommendation (cached)
    st.markdown("### 🤖 GPT Recommendation")
    st.success(st.session_state["gpt_recommendation"])

    # Optional: show the full stats table
    if st.button("📊 Show Detailed Stats Table"):
        st.dataframe(summary_df)

    # User selection
    method = st.radio("Choose Effectiveness Method", ["AM", "GM"])
    if st.button("Next"):
        st.session_state.method = method
        st.session_state.step += 1
        st.rerun()

# ----------------------------
# Step 3: Effectiveness Score
# ----------------------------
elif st.session_state.step == 3:
    st.header("Step 3: Calculate Effectiveness Score")
    df = st.session_state.df
    method = st.session_state.method
    effectiveness_df = pd.DataFrame()
    for msg in sorted(set(col.split("_")[0] for col in df.columns if col.startswith("M"))):
        cols = [f"{msg}_Differentiated", f"{msg}_Believable", f"{msg}_Motivating"]
        if method == "AM":
            effectiveness_df[f"{msg}_Effectiveness"] = df[cols].mean(axis=1)
        else:
            effectiveness_df[f"{msg}_Effectiveness"] = df[cols].replace(0, 0.01).prod(axis=1)**(1/3)
    st.session_state.effectiveness_df = effectiveness_df

    # --- Visualization ---
    st.subheader("📊 Average Effectiveness by Message")
    avg_scores = effectiveness_df.mean().reset_index()
    avg_scores.columns = ["Message", "Average Effectiveness"]
    avg_scores["Message"] = avg_scores["Message"].str.replace("_Effectiveness", "", regex=False)
    # Sort in descending order
    avg_scores = avg_scores.sort_values(by="Average Effectiveness", ascending=False)
    sorted_messages = avg_scores["Message"].tolist()  # Keep sorted order for plotting

    # Plot
    fig, ax = plt.subplots(figsize=(10, 4))
    sns.barplot(data=avg_scores, x="Message", y="Average Effectiveness", palette="viridis", ax=ax, order=sorted_messages)
    # ✅ Correct label placement using bar patches
    for patch in ax.patches:
        height = patch.get_height()
        ax.text(
            patch.get_x() + patch.get_width() / 2,
            height + 0.02,
            f"{height:.2f}",
            ha='center',
            va='bottom',
            fontsize=9,
            color='black'
        )
    ax.set_title("Average Effectiveness Scores (Sorted)")
    ax.set_ylabel("Score")
    ax.set_xlabel("Message")
    plt.xticks(rotation=45)
    st.pyplot(fig)

    if st.button("Next"):
        st.session_state.step += 1
        st.rerun()

# ----------------------------
# Step 4: Flatliner Removal
# ----------------------------
elif st.session_state.step == 4:
    st.header("Step 4: Remove Flatliners?")

    # Always work from original unfiltered effectiveness_df
    original_df = st.session_state.get("original_effectiveness_df")
    if original_df is None:
        original_df = st.session_state.effectiveness_df.copy()
        st.session_state.original_effectiveness_df = original_df

    option = st.radio("Remove respondents with low variance?", ["No", "Yes"], index=0)
    threshold = st.number_input("Variance Threshold", min_value=0.0, max_value=10.0, value=0.05, step=0.01)

    if option == "Yes":
        # Calculate from original each time
        row_variances = original_df.var(axis=1, ddof=0)
        retained_df = original_df[row_variances >= threshold]
        removed_count = original_df.shape[0] - retained_df.shape[0]
        st.write(f"✅ Retained: {retained_df.shape[0]} rows")
        st.write(f"🗑 Removed: {removed_count} rows")
    else:
        st.info(f"Retained all {original_df.shape[0]} respondents (no flatliners removed).")

    # Only apply filtering on button click
    if st.button("Next"):
        if option == "Yes":
            st.session_state.effectiveness_df = retained_df
        else:
            st.session_state.effectiveness_df = original_df
        # Cleanup original store
        del st.session_state.original_effectiveness_df
        st.session_state.step += 1
        st.rerun()

# ----------------------------
# Step 5: Binarization with GPT Recommendations
# ----------------------------
elif st.session_state.step == 5:
    st.header("Step 5: Binarization")
    df = st.session_state.effectiveness_df.copy()
    st.write(f"🧮 Total respondents: {df.shape[0]}")

    # --- GPT Recommendation for binarization method choice ---
    if "binarization_gpt_recommendation" not in st.session_state:
        try:
            gpt_prompt = ("You are a data analytics expert specializing in survey analysis. "
                          "For TURF analysis of PET message effectiveness scores, recommend which binarization method would be most appropriate among these options:\n\n"
                          "1) T2B (Top 2 Box) - Simple cutoff at score > 5\n"
                          "2) Index method - Uses percentage threshold above each respondent's personal mean\n"
                          "3) Segment Based Index (GMM) - Uses Gaussian Mixture Model clustering with segment-specific thresholds\n\n"
                          "Explain the criteria for choosing each method and scenarios where each is best suited. "
                          "Provide your response in 3 concise bullet points, one for each method.")
            
            key = st.secrets["openai_key"]
            client = openai.OpenAI(api_key=key)
            response = client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are an expert in data analytics and survey data analysis."},
                    {"role": "user", "content": gpt_prompt}
                ],
                temperature=0.3
            )
            st.session_state["binarization_gpt_recommendation"] = response.choices[0].message.content
        except Exception as e:
            st.session_state["binarization_gpt_recommendation"] = f"⚠️ GPT recommendation failed: {e}"

    method = st.radio("Choose binarization method", ["T2B", "Index (X% above mean)", "Segment Based Index (GMM)"])

    # Method-specific parameters with GPT recommendations
    if method == "T2B":
        st.info("**T2B Method:** Converts effectiveness scores above 5 to 1, else 0. Simple and interpretable for Likert scale data.")
        
    elif method == "Index (X% above mean)":
        percent_above = st.number_input("Set Index Threshold (% above respondent mean)", min_value=0.0, max_value=100.0, value=5.0, step=0.5)

        # Concise GPT recommendation for Index threshold (0-15% range)
        if "index_threshold_gpt" not in st.session_state:
            try:
                gpt_index_prompt = ("For pharmaceutical message testing using Index binarization, recommend a specific threshold percentage between 0-15% above each respondent's personal mean. "
                                    "Provide one specific percentage with brief rationale in 1-2 sentences.")
                
                key = st.secrets["openai_key"]
                client = openai.OpenAI(api_key=key)
                response_index = client.chat.completions.create(
                    model="gpt-4",
                    messages=[
                        {"role": "system", "content": "You are an expert in pharmaceutical survey analytics. Always recommend thresholds between 0-15%. Be concise."},
                        {"role": "user", "content": gpt_index_prompt}
                    ],
                    temperature=0.3
                )
                st.session_state["index_threshold_gpt"] = response_index.choices[0].message.content
            except Exception as e:
                st.session_state["index_threshold_gpt"] = f"⚠️ GPT recommendation failed: {e}"

        st.markdown("#### 🤖 GPT Recommended Index Threshold")
        st.info(st.session_state["index_threshold_gpt"])

    elif method == "Segment Based Index (GMM)":
        threshold_pct = st.number_input("Set threshold for segment index binarization (% above segment mean)", min_value=0.0, max_value=100.0, value=5.0, step=0.5)
        n_clusters = st.number_input("Number of clusters (segments)", min_value=2, max_value=10, value=4, step=1)

        # Concise GPT recommendation for segment-based parameters (0-15% threshold range)
        if "segment_params_gpt" not in st.session_state:
            try:
                gpt_segment_prompt = ("For pharmaceutical Segment Based Index binarization using GMM clustering, "
                                    "recommend: 1) threshold percentage (0-15% above segment mean), and 2) number of clusters (2-10). "
                                    "Provide specific recommendations in 2 brief sentences.")
                
                key = st.secrets["openai_key"]
                client = openai.OpenAI(api_key=key)
                response_segment = client.chat.completions.create(
                    model="gpt-4",
                    messages=[
                        {"role": "system", "content": "You are an expert in pharmaceutical survey analytics. Keep threshold recommendations between 0-15%. Be concise."},
                        {"role": "user", "content": gpt_segment_prompt}
                    ],
                    temperature=0.3
                )
                st.session_state["segment_params_gpt"] = response_segment.choices[0].message.content
            except Exception as e:
                st.session_state["segment_params_gpt"] = f"⚠️ GPT recommendation failed: {e}"

        st.markdown("#### 🤖 GPT Recommended Segment Based Index Parameters")
        st.info(st.session_state["segment_params_gpt"])

    # Execute binarization based on selected method
    if method == "T2B":
        binarized_df = df.applymap(lambda x: 1 if x > 5 else 0)
    elif method == "Index (X% above mean)":
        multiplier = 1 + percent_above / 100
        binarized_df = df.copy()
        for idx, row in df.iterrows():
            threshold = row.mean() * multiplier
            binarized_df.loc[idx] = row.apply(lambda x: 1 if x >= threshold else 0)
    elif method == "Segment Based Index (GMM)":
        from sklearn.mixture import GaussianMixture
        st.markdown("📌 Segmenting respondents using GMM Clustering")
        
        threshold_index = 100 + threshold_pct
        gmm = GaussianMixture(n_components=n_clusters, random_state=42)
        segments = gmm.fit_predict(df)
        df["Segment"] = segments
        
        binarized_df = pd.DataFrame(0, index=df.index, columns=df.columns.drop("Segment"))
        for msg in binarized_df.columns:
            segment_means = df.groupby("Segment")[msg].mean()
            for idx in df.index:
                seg = df.loc[idx, "Segment"]
                seg_mean = segment_means.loc[seg]
                if seg_mean == 0:
                    index_score = 0
                else:
                    index_score = (df.loc[idx, msg] / seg_mean) * 100
                if index_score > threshold_index:
                    binarized_df.loc[idx, msg] = 1
        df.drop(columns=["Segment"], inplace=True)

    st.session_state.binarized_df = binarized_df

    # --- Visualization ---
    st.subheader("📊 Message-wise Reach (% of 1s after binarization)")
    percentage_ones = (binarized_df.sum(axis=0) / binarized_df.shape[0] * 100).reset_index()
    percentage_ones.columns = ["Message", "Percentage"]
    percentage_ones["Message"] = percentage_ones["Message"].str.replace("_Effectiveness", "", regex=False)
    percentage_ones = percentage_ones.sort_values(by="Percentage", ascending=False)

    fig, ax = plt.subplots(figsize=(10, 4))
    sns.barplot(data=percentage_ones, x="Message", y="Percentage", palette="Blues_d", ax=ax)
    for patch in ax.patches:
        height = patch.get_height()
        ax.text(
            patch.get_x() + patch.get_width() / 2,
            height + 1,
            f"{height:.1f}%",
            ha='center',
            va='bottom',
            fontsize=9
        )
    ax.set_title("Binarized Reach per Message")
    ax.set_ylabel("Percentage of 1s")
    ax.set_xlabel("Message")
    plt.xticks(rotation=45)
    st.pyplot(fig)

    if st.button("Next"):
        st.session_state.step += 1
        st.rerun()

    # General GPT guidance moved to bottom
    st.markdown("### 🤖 General Binarization Method Guide")
    st.success(st.session_state.get("binarization_gpt_recommendation", "Loading guidance..."))

# ----------------------------
# Step 6: TURF Analysis
# ----------------------------
elif st.session_state.step == 6:
    st.header("Step 6: TURF Analysis (Reach by Bundle Size)")
    df = st.session_state.binarized_df
    message_reach = df.sum() / df.shape[0]

    turf_summary_rows = []
    best_combos = {}

    for k in range(1, 6):
        # Greedy TURF algorithm
        greedy_combo = []
        remaining = list(df.columns)
        turf_input_df = df.copy()

        for _ in range(k):
            best_var = None
            best_reach = -1
            best_freq = -1

            for var in remaining:
                test_combo = greedy_combo + [var]
                test_df = turf_input_df[test_combo]
                respondent_reach = (test_df.sum(axis=1) > 0)
                reach = respondent_reach.mean()

                # ✅ Avg frequency among reached respondents
                if respondent_reach.sum() > 0:
                    freq = test_df[respondent_reach].sum(axis=1).mean()
                else:
                    freq = 0

                if reach > best_reach or (reach == best_reach and freq > best_freq):
                    best_var = var
                    best_reach = reach
                    best_freq = freq

            greedy_combo.append(best_var)
            remaining.remove(best_var)

        # Final best combo for this k
        final_df = turf_input_df[greedy_combo]
        respondent_reach = (final_df.sum(axis=1) > 0)
        reach_pct = round(respondent_reach.mean() * 100, 2)
        avg_freq = round(final_df[respondent_reach].sum(axis=1).mean(), 2)

        # Sort combo by message-level reach (for display only)
        sorted_combo = sorted(greedy_combo, key=lambda m: -message_reach[m])
        cleaned_combo = ", ".join([m.split("_")[0] for m in sorted_combo])

        turf_summary_rows.append({
            "Messages in Bundle": k,
            "Reach (%)": reach_pct,
            "Avg Frequency": avg_freq,
            "Best Combination": cleaned_combo
        })

        best_combos[k] = sorted_combo

    # Final summary table
    turf_summary = pd.DataFrame(turf_summary_rows)

    # Save to session
    st.session_state.turf_summary = turf_summary
    st.session_state.best_combos = best_combos

    # Display summary
    st.subheader("📊 TURF Reach Summary")
    st.dataframe(turf_summary)

    # Plot Reach Curve
    sns.lineplot(data=turf_summary, x="Messages in Bundle", y="Reach (%)", marker="o", color="green")
    plt.title("Reach by Bundle Size")
    plt.ylabel("Reach (%)")
    plt.xlabel("Messages in Bundle")
    st.pyplot(plt.gcf())

    # --- 🤖 GPT Recommendation ---
    try:
        prompt = "You are a pharma messaging strategy expert. Below is a TURF analysis output showing reach by number of messages in a bundle. Based on this, suggest the optimal number of messages (one single number) that balances reach and message overload. Justify in 2-3 sentences why this number is optimal.\n\n"
        
        for _, row in turf_summary.iterrows():
            prompt += f"{int(row['Messages in Bundle'])} messages → Reach: {row['Reach (%)']}%, Avg Freq: {row['Avg Frequency']}, Best Combo: {row['Best Combination']}\n"
        
        prompt += "\nPlease recommend ONE optimal number of messages to proceed with."

        key = st.secrets["openai_key"]
        client = openai.OpenAI(api_key=key)
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a marketing analytics and messaging expert."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3
        )
        
        gpt_recommendation = response.choices[0].message.content
        st.markdown("### 🤖 GPT Recommendation")
        st.success(gpt_recommendation)
        st.session_state["gpt_recommendation"] = gpt_recommendation
        
    except Exception as e:
        st.warning(f"⚠️ GPT recommendation skipped: {e}")

    if st.button("Next"):
        st.session_state.step += 1
        st.rerun()

# ----------------------------
# Step 7: Monte Carlo Simulation
# ----------------------------
elif st.session_state.step == 7:
    st.header("Step 7: Monte Carlo Simulation")
    
    run_sim = st.radio("Run simulation?", ["Yes", "No"])

    if run_sim == "Yes":
        bundle = st.slider("Bundle size", 1, 5, 3)
        iterations = st.slider("Iterations", 10, 100, 25)
        
        df = st.session_state.binarized_df
        wins = []
        
        for i in range(iterations):
            sample = df.sample(frac=0.8, replace=False, random_state=i)
            combos = list(itertools.combinations(df.columns, bundle))
            best = max(combos, key=lambda c: (sample[list(c)].sum(axis=1) > 0).mean())
            wins.append(tuple(sorted(best)))
        
        counts = Counter(wins).most_common()
        
        st.subheader("Top Monte Carlo Combos")
        for combo, freq in counts[:5]:
            st.write(f"{', '.join([m.split('_')[0] for m in combo])} → {freq} wins")
        
        # ✅ Check if TURF combo matches simulation top combos
        original = tuple(sorted(st.session_state.best_combos[bundle]))
        if original in [c for c, _ in counts]:
            st.success("✅ Match with TURF result — stable")
            st.session_state["monte_carlo_result"] = "✅ Match with TURF result — stable"
        else:
            st.warning("⚠️ No match — result may not be stable")
            st.session_state["monte_carlo_result"] = "⚠️ No match — result may not be stable"
    else:
        st.info("Simulation skipped.")
        st.session_state["monte_carlo_result"] = "Simulation skipped."

    if st.button("Next"):
        st.session_state.step += 1
        st.rerun()

# ----------------------------
# Step 8: Final Summary
# ----------------------------
elif st.session_state.step == 8:
    st.header("✅ Final Summary")
    
    turf_summary = st.session_state.turf_summary
    best_combos = st.session_state.best_combos
    monte_carlo_confidence = st.session_state.get("monte_carlo_result", "Monte Carlo confidence not available.")

    # --- Display TURF Table and Combos ---
    st.subheader("TURF Results")
    st.dataframe(turf_summary)

    st.subheader("Best Combos by Bundle Size")
    for k, combo in best_combos.items():
        st.markdown(f"- **{k} messages** → {', '.join([m.split('_')[0] for m in combo])}")

    # --- STEP 1: Extract bundle size from GPT step 6 response ---
    import re
    raw_gpt_text = st.session_state.get("gpt_recommendation", "")
    bundle_size_match = re.search(r'\b(\d)\s*messages?\b', raw_gpt_text)
    
    if bundle_size_match:
        bundle_size = int(bundle_size_match.group(1))
    else:
        bundle_size = int(turf_summary.sort_values("Reach (%)", ascending=False).iloc[0]["Messages in Bundle"])

    # --- STEP 2: Fetch best combo and reach for that size ---
    row = turf_summary[turf_summary["Messages in Bundle"] == bundle_size].iloc[0]
    best_combo = best_combos[bundle_size]
    best_combo_labels = ", ".join([m.split("_")[0] for m in best_combo])

    # --- STEP 3: Build GPT Prompt with selected bundle size ---
    gpt_prompt = (
        "You are a pharma marketing strategist. Based on the TURF analysis and Monte Carlo simulation, "
        "summarize the recommended bundle strategy as three concise bullet points:\n\n"
        f"- TURF output: Best bundle size = {bundle_size}, Reach = {row['Reach (%)']}%\n"
        f"- Recommended message sequence: {best_combo_labels}\n"
        f"- Monte Carlo simulation confidence: {monte_carlo_confidence}\n\n"
        "👉 Please respond in exactly 3 clear, standalone bullet points covering:\n"
        "1. The best bundle size and why\n"
        "2. The message combination in recommended sequence\n"
        "3. Confidence in recommendation based on Monte Carlo simulation\n\n"
        "Format:\n"
        "• [Bundle Size] ...\n"
        "• [Message Sequence] ...\n"
        "• [Confidence] ..."
    )

    try:
        key = st.secrets["openai_key"]
        client = openai.OpenAI(api_key=key)
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an expert in healthcare message optimization."},
                {"role": "user", "content": gpt_prompt}
            ],
            temperature=0.3
        )
        
        gpt_final_bullets = response.choices[0].message.content.strip()
        st.session_state["gpt_final_summary"] = gpt_final_bullets
        
    except Exception as e:
        gpt_final_bullets = f"⚠️ GPT recommendation failed: {e}"
        st.session_state["gpt_final_summary"] = gpt_final_bullets

    # --- STEP 4: Display formatted bullets ---
    st.subheader("🤖 Final GPT Recommendation")
    bullets = [line.strip("•").strip("-").strip() for line in gpt_final_bullets.split("\n") if line.strip()]
    for bullet in bullets:
        st.markdown(f"- {bullet}")

    # --- STEP 5: Create chart for PPT ---
    fig, ax = plt.subplots(figsize=(6, 3))
    ax.plot(turf_summary["Messages in Bundle"], turf_summary["Reach (%)"], marker='o', color='green')
    ax.set_title("TURF Reach by Bundle Size")
    ax.set_xlabel("Messages in Bundle")
    ax.set_ylabel("Reach (%)")
    plt.tight_layout()
    
    chart_buf = io.BytesIO()
    plt.savefig(chart_buf, format='png')
    chart_buf.seek(0)
    plt.close()

    # --- STEP 6: Build PPT ---
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TURF Analysis Summary"
    slide.placeholders[1].text = "Auto-generated summary of reach and GPT recommendation."

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5)).text = "TURF Reach Curve"
    slide.shapes.add_picture(chart_buf, Inches(1), Inches(1.2), width=Inches(6.5))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(5)).text_frame
    tf.text = "Best Message Combinations:\n"
    for _, row in turf_summary.iterrows():
        tf.add_paragraph().text = f"{int(row['Messages in Bundle'])} messages → {row['Best Combination']}"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(4)).text_frame
    tf2.text = "🤖 GPT Final Recommendation:\n"
    for bullet in bullets:
        tf2.add_paragraph().text = bullet

    ppt_buf = io.BytesIO()
    prs.save(ppt_buf)
    ppt_buf.seek(0)

    st.download_button(
        label="📥 Download TURF Summary PPT",
        data=ppt_buf,
        file_name="TURF_Summary_Presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    if st.button("🔁 Restart"):
        for k in list(st.session_state.keys()):
            del st.session_state[k]
        st.rerun()
